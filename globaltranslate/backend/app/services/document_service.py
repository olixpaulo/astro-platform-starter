"""Extração de texto de documentos (PDF, DOCX, TXT, PPTX) e tradução em blocos."""
import io

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from app.services.translation_service import translate_text

SUPPORTED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "text/plain": "txt",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

# Limite por chamada ao modelo para não exceder a janela de contexto
CHUNK_SIZE = 6_000


def extract_text(file_bytes: bytes, content_type: str) -> str:
    kind = SUPPORTED_TYPES.get(content_type)
    if kind == "pdf":
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    if kind == "docx":
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)
    if kind == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)
    if kind == "txt":
        return file_bytes.decode("utf-8", errors="replace")
    raise ValueError(f"Tipo de documento não suportado: {content_type}")


def _chunk(text: str, size: int = CHUNK_SIZE) -> list[str]:
    """Divide o texto em blocos respeitando quebras de parágrafo."""
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0
    for paragraph in text.split("\n\n"):
        if current_len + len(paragraph) > size and current:
            chunks.append("\n\n".join(current))
            current, current_len = [], 0
        current.append(paragraph)
        current_len += len(paragraph) + 2
    if current:
        chunks.append("\n\n".join(current))
    return chunks


async def translate_document_text(text: str, source_lang: str, target_lang: str) -> str:
    translated: list[str] = []
    for chunk in _chunk(text):
        if not chunk.strip():
            translated.append(chunk)
            continue
        result = await translate_text(chunk, source_lang, target_lang)
        translated.append(result["translated_text"])
    return "\n\n".join(translated)
